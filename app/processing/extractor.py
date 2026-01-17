# app/processing/extractor.py
"""Document text extraction for various file types"""

import io
from typing import Optional, Dict, Any, List
from dataclasses import dataclass
import logging

# PDF
import fitz  # PyMuPDF

# Office documents
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# Other formats
import pandas as pd
from bs4 import BeautifulSoup
import json
import csv

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    """Result of text extraction"""
    text: str
    metadata: Dict[str, Any]
    page_count: Optional[int] = None
    tables: Optional[List[Dict]] = None
    error: Optional[str] = None


class DocumentExtractor:
    """
    Extract text content from various document types.
    """
    
    # MIME type to handler mapping
    HANDLERS = {
        'application/pdf': 'extract_pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'extract_docx',
        'application/msword': 'extract_docx',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'extract_xlsx',
        'application/vnd.ms-excel': 'extract_xlsx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'extract_pptx',
        'application/vnd.ms-powerpoint': 'extract_pptx',
        'text/plain': 'extract_text',
        'text/csv': 'extract_csv',
        'text/markdown': 'extract_text',
        'text/html': 'extract_html',
        'application/json': 'extract_json',
        'image/png': 'extract_image',
        'image/jpeg': 'extract_image',
        'image/jpg': 'extract_image',
    }
    
    def __init__(self, ocr_enabled: bool = True):
        self.ocr_enabled = ocr_enabled
    
    async def extract(
        self,
        content: bytes,
        mime_type: str,
        filename: str
    ) -> ExtractionResult:
        """
        Extract text from document.
        
        Args:
            content: Raw file bytes
            mime_type: MIME type of the file
            filename: Original filename
            
        Returns:
            ExtractionResult with extracted text and metadata
        """
        # Find handler
        handler_name = self.HANDLERS.get(mime_type)
        
        # Fallback: detect by extension
        if not handler_name:
            ext = '.' + filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
            ext_mapping = {
                '.pdf': 'extract_pdf',
                '.docx': 'extract_docx',
                '.doc': 'extract_docx',
                '.xlsx': 'extract_xlsx',
                '.xls': 'extract_xlsx',
                '.pptx': 'extract_pptx',
                '.ppt': 'extract_pptx',
                '.txt': 'extract_text',
                '.md': 'extract_text',
                '.csv': 'extract_csv',
                '.tsv': 'extract_csv',
                '.html': 'extract_html',
                '.htm': 'extract_html',
                '.json': 'extract_json',
                '.png': 'extract_image',
                '.jpg': 'extract_image',
                '.jpeg': 'extract_image',
            }
            handler_name = ext_mapping.get(ext)
        
        if not handler_name:
            return ExtractionResult(
                text="",
                metadata={"error": f"Unsupported file type: {mime_type}"},
                error=f"Unsupported file type: {mime_type}"
            )
        
        try:
            handler = getattr(self, handler_name)
            return await handler(content, filename)
        except Exception as e:
            logger.error(f"Extraction error for {filename}: {e}")
            return ExtractionResult(
                text="",
                metadata={"error": str(e)},
                error=str(e)
            )
    
    async def extract_pdf(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from PDF"""
        doc = fitz.open(stream=content, filetype="pdf")

        text_parts = []
        page_texts = {}
        page_count = len(doc)  # Store before closing

        for page_num in range(page_count):
            page = doc[page_num]
            page_text = page.get_text("text")

            # Try OCR if no text extracted and OCR is enabled
            if not page_text.strip() and self.ocr_enabled:
                try:
                    page_text = self._ocr_page(page)
                except Exception as e:
                    logger.warning(f"OCR failed for page {page_num + 1}: {e}")

            if page_text.strip():
                text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                page_texts[page_num + 1] = page_text

        # Extract metadata before closing
        metadata = {
            "title": doc.metadata.get("title"),
            "author": doc.metadata.get("author"),
            "subject": doc.metadata.get("subject"),
            "creator": doc.metadata.get("creator"),
            "page_count": page_count,
            "page_texts": page_texts
        }

        doc.close()

        return ExtractionResult(
            text="\n\n".join(text_parts),
            metadata=metadata,
            page_count=page_count
        )
    
    def _ocr_page(self, page) -> str:
        """OCR a PDF page using Tesseract"""
        try:
            import pytesseract
            from PIL import Image
            
            # Render page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # OCR with Thai + English
            text = pytesseract.image_to_string(img, lang='eng+tha')
            return text
        except ImportError:
            logger.warning("pytesseract not available for OCR")
            return ""
    
    async def extract_docx(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from Word document"""
        doc = DocxDocument(io.BytesIO(content))
        
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Check if it's a heading
                if para.style and para.style.name.startswith('Heading'):
                    text_parts.append(f"\n## {para.text}\n")
                else:
                    text_parts.append(para.text)
        
        # Extract tables
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                tables.append(table_data)
                # Add table as text
                table_text = self._table_to_text(table_data)
                text_parts.append(f"\n[Table]\n{table_text}\n")
        
        # Metadata
        core_props = doc.core_properties
        metadata = {
            "title": core_props.title,
            "author": core_props.author,
            "subject": core_props.subject,
            "created": str(core_props.created) if core_props.created else None,
            "modified": str(core_props.modified) if core_props.modified else None,
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(tables)
        }
        
        return ExtractionResult(
            text="\n".join(text_parts),
            metadata=metadata,
            tables=tables
        )
    
    async def extract_xlsx(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from Excel spreadsheet"""
        wb = load_workbook(io.BytesIO(content), data_only=True)
        
        text_parts = []
        tables = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # Get all data
            data = []
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(row_data):  # Skip empty rows
                    data.append(row_data)
            
            if data:
                tables.append({"sheet": sheet_name, "data": data})
                
                # Convert to text
                text_parts.append(f"\n## Sheet: {sheet_name}\n")
                text_parts.append(self._table_to_text(data))
        
        wb.close()
        
        return ExtractionResult(
            text="\n".join(text_parts),
            metadata={
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames
            },
            tables=tables
        )
    
    async def extract_pptx(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from PowerPoint presentation"""
        prs = Presentation(io.BytesIO(content))
        
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract from tables
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    slide_text.append(self._table_to_text(table_data))
            
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        
        return ExtractionResult(
            text="\n\n".join(text_parts),
            metadata={
                "slide_count": len(prs.slides)
            },
            page_count=len(prs.slides)
        )
    
    async def extract_text(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from plain text file"""
        # Try different encodings
        encodings = ['utf-8', 'utf-8-sig', 'cp1252', 'iso-8859-1', 'tis-620']
        
        text = None
        for encoding in encodings:
            try:
                text = content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if text is None:
            text = content.decode('utf-8', errors='replace')
        
        return ExtractionResult(
            text=text,
            metadata={"encoding": encoding if text else "unknown"}
        )
    
    async def extract_csv(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from CSV file"""
        # Detect delimiter
        sample = content[:4096].decode('utf-8', errors='replace')
        delimiter = ',' if sample.count(',') > sample.count('\t') else '\t'
        
        try:
            df = pd.read_csv(io.BytesIO(content), delimiter=delimiter)
            
            # Convert to text representation
            text_parts = [f"Columns: {', '.join(df.columns)}"]
            text_parts.append(f"Rows: {len(df)}")
            text_parts.append("\nData:\n")
            text_parts.append(df.to_string(index=False, max_rows=1000))
            
            return ExtractionResult(
                text="\n".join(text_parts),
                metadata={
                    "columns": list(df.columns),
                    "row_count": len(df),
                    "column_count": len(df.columns)
                },
                tables=[{"data": df.values.tolist(), "columns": list(df.columns)}]
            )
        except Exception as e:
            # Fallback to plain text
            return await self.extract_text(content, filename)
    
    async def extract_html(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from HTML file"""
        text = content.decode('utf-8', errors='replace')
        soup = BeautifulSoup(text, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        extracted_text = soup.get_text(separator='\n')
        
        # Clean up whitespace
        lines = [line.strip() for line in extracted_text.splitlines()]
        extracted_text = '\n'.join(line for line in lines if line)
        
        # Get title
        title = soup.title.string if soup.title else None
        
        return ExtractionResult(
            text=extracted_text,
            metadata={
                "title": title,
                "has_tables": len(soup.find_all('table')) > 0
            }
        )
    
    async def extract_json(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from JSON file"""
        try:
            data = json.loads(content.decode('utf-8'))
            
            # Pretty print JSON
            text = json.dumps(data, indent=2, ensure_ascii=False)
            
            return ExtractionResult(
                text=text,
                metadata={
                    "type": type(data).__name__,
                    "keys": list(data.keys()) if isinstance(data, dict) else None,
                    "length": len(data) if isinstance(data, (list, dict)) else None
                }
            )
        except json.JSONDecodeError as e:
            return ExtractionResult(
                text=content.decode('utf-8', errors='replace'),
                metadata={"error": f"Invalid JSON: {e}"},
                error=f"Invalid JSON: {e}"
            )
    
    async def extract_image(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from Image using OCR"""
        if not self.ocr_enabled:
            return ExtractionResult(
                text="",
                metadata={"error": "OCR disabled"},
                error="OCR disabled"
            )
            
        try:
            import pytesseract
            from PIL import Image
            
            # Load image
            image = Image.open(io.BytesIO(content))
            
            # OCR with Thai + English
            text = pytesseract.image_to_string(image, lang='eng+tha')
            
            return ExtractionResult(
                text=text,
                metadata={
                    "format": image.format,
                    "mode": image.mode,
                    "size": image.size,
                    "width": image.width,
                    "height": image.height
                }
            )
        except ImportError:
            return ExtractionResult(
                text="",
                metadata={"error": "pytesseract not available for OCR"},
                error="pytesseract not available"
            )
        except Exception as e:
            logger.error(f"OCR error for {filename}: {e}")
            return ExtractionResult(
                text="",
                metadata={"error": str(e)},
                error=str(e)
            )
    
    def _table_to_text(self, data: List[List[str]]) -> str:
        """Convert table data to text representation"""
        if not data:
            return ""
        
        # Use first row as header
        lines = [" | ".join(str(cell) for cell in data[0])]
        lines.append("-" * len(lines[0]))
        
        for row in data[1:]:
            lines.append(" | ".join(str(cell) for cell in row))
        
        return "\n".join(lines)
